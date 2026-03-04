"""
Tool for sewing images into a PPT file.
Originally by Hylz – rewritten for web-based CherryPicker.

Bug fixes vs original:
- Guard against empty img_paths (no crops generated yet)
- Guard against unreadable images when computing dimensions
- Use sorted(seqs.keys()) for deterministic ordering
- Avoid crash when placeholder image is missing
- Better logging instead of bare print
"""

import glob
import logging
import os

import cv2
from pptx import Presentation
from pptx.util import Mm

logger = logging.getLogger("cherrypicker.ppt")


def make_ppt(config: dict) -> None:
    root_path = config["output_crop_path"]
    if not os.path.isdir(root_path):
        raise FileNotFoundError(
            f"Crop output directory not found: {root_path}. Run 'Make All Crops' first."
        )

    method_names = [m["name"] for m in config["methods"]]
    method_cnt = len(method_names)

    SMALL_CNT = config.get("small_cnt", 3)
    GROUPS_PER_PAGE = config.get("groups_per_page", 5)
    SLIDE_W = config.get("slide_w", 210)
    SLIDE_H = config.get("slide_h", 297)

    GROUP_HORI_GAP_RATIO = config.get("group_hori_gap_ratio", 0.05)
    GROUP_VERT_GAP_RATIO = config.get("group_vert_gap_ratio", 0.07)
    SMALL_HORI_GAP_RATIO = config.get("small_hori_gap_ratio", 0.05)
    SMALL_VERT_GAP_RATIO = config.get("small_vert_gap_ratio", 0.05)

    AREA_W = SLIDE_W * 0.8
    AREA_H = SLIDE_H * 0.95  # FIX: was SLIDE_W * 0.95 (typo)

    placeholder = config.get("placeholder_path", "placeholder.png")

    # ---- Collect image paths per method ---------------------------------
    img_paths: list[list[str]] = []

    for m in method_names:
        method_dir = os.path.join(root_path, m)
        if not os.path.isdir(method_dir):
            logger.warning("Method directory missing: %s", method_dir)
            img_paths.append([])
            continue

        all_imgs = sorted(
            glob.glob(os.path.join(method_dir, "*.png"))
            + glob.glob(os.path.join(method_dir, "*.jpg"))
        )

        seqs: dict[str, dict] = {}
        for img in all_imgs:
            base = os.path.basename(img)
            seq = base.split("_")[0]
            if seq not in seqs:
                seqs[seq] = {"big": None, "small": []}
            if "crop" in base:
                if len(seqs[seq]["small"]) >= SMALL_CNT:
                    logger.warning(
                        "Sequence %s has more than %d crops – extras ignored.", seq, SMALL_CNT
                    )
                else:
                    seqs[seq]["small"].append(img)
            elif "full" in base:
                seqs[seq]["big"] = img

        my_imgs: list[str] = []
        for seq in sorted(seqs.keys()):
            big = seqs[seq]["big"]
            if big is None:
                logger.warning("No full image for sequence %s – skipping.", seq)
                continue
            my_imgs.append(big)
            my_imgs.extend(seqs[seq]["small"])
            # Pad with placeholder if not enough crops
            for _ in range(SMALL_CNT - len(seqs[seq]["small"])):
                my_imgs.append(placeholder)
        img_paths.append(my_imgs)

    # Validate
    if not img_paths or not img_paths[0]:
        raise ValueError("No cropped images found. Run 'Make All Crops' first.")

    stride = SMALL_CNT + 1
    group_cnt = len(img_paths[0]) // stride
    for pthlist in img_paths:
        if len(pthlist) != group_cnt * stride:
            raise ValueError(
                f"Image count mismatch: expected {group_cnt * stride}, got {len(pthlist)}"
            )

    page_cnt = (group_cnt + GROUPS_PER_PAGE - 1) // GROUPS_PER_PAGE

    # ---- Compute layout dimensions from first images --------------------
    big_sample = cv2.imread(img_paths[0][0])
    if big_sample is None:
        raise FileNotFoundError(f"Cannot read big image: {img_paths[0][0]}")
    bighpix, bigwpix = big_sample.shape[:2]

    small_sample = cv2.imread(img_paths[0][1])
    if small_sample is None:
        raise FileNotFoundError(f"Cannot read small image: {img_paths[0][1]}")
    smallhpix, smallwpix = small_sample.shape[:2]

    group_w = AREA_W / GROUPS_PER_PAGE
    bigw = group_w / (1 + GROUP_HORI_GAP_RATIO)
    bigh = bigw * bighpix / max(bigwpix, 1)
    smallw = bigw / (SMALL_CNT + (SMALL_CNT - 1) * SMALL_HORI_GAP_RATIO) if SMALL_CNT else bigw
    smallh = smallw * smallhpix / max(smallwpix, 1)
    group_h = (bigh + smallh * (1 + SMALL_VERT_GAP_RATIO)) * (1 + GROUP_VERT_GAP_RATIO)

    # ---- Build the presentation -----------------------------------------
    prs = Presentation()
    prs.slide_width = Mm(SLIDE_W)
    prs.slide_height = Mm(SLIDE_H)
    blank_layout = prs.slide_layouts[6]

    for pagenum in range(page_cnt):
        slide = prs.slides.add_slide(blank_layout)

        for row_i in range(method_cnt):
            big_top = (SLIDE_H - AREA_H) / 2 + group_h * row_i
            small_top = big_top + bigh + smallh * SMALL_VERT_GAP_RATIO

            # Method name label
            text_left = 0
            text_top = (big_top + small_top) / 2
            text_width = (SLIDE_W - AREA_W) / 2
            text_height = small_top - big_top
            txbox = slide.shapes.add_textbox(
                Mm(text_left), Mm(text_top), Mm(text_width), Mm(text_height)
            )
            p = txbox.text_frame.add_paragraph()
            p.text = method_names[row_i]

            start = pagenum * GROUPS_PER_PAGE * stride
            end = min(start + GROUPS_PER_PAGE * stride, len(img_paths[row_i]))
            imglist = img_paths[row_i][start:end]
            page_groupcnt = len(imglist) // stride

            for group_i in range(page_groupcnt):
                big_left = (SLIDE_W - AREA_W) / 2 + group_w * group_i
                big_path = imglist[group_i * stride]
                if os.path.isfile(big_path):
                    slide.shapes.add_picture(
                        big_path, Mm(big_left), Mm(big_top),
                        width=Mm(bigw), height=Mm(bigh),
                    )
                for small_i in range(SMALL_CNT):
                    small_left = big_left + (1 + SMALL_HORI_GAP_RATIO) * smallw * small_i
                    small_path = imglist[group_i * stride + small_i + 1]
                    if os.path.isfile(small_path):
                        slide.shapes.add_picture(
                            small_path, Mm(small_left), Mm(small_top),
                            width=Mm(smallw), height=Mm(smallh),
                        )

    output_path = config["output_ppt_path"]
    prs.save(output_path)
    logger.info("PPT saved to %s", output_path)